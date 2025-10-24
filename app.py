# app/app.py
# ---------------------------------------------------------
# TechTitans – Renewable Valuation Dashboard (P75-ready)
# Avangrid Hackathon 2025
# ---------------------------------------------------------

import io
from pathlib import Path
import numpy as np
import pandas as pd
import plotly.express as px
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt

# ---------------------------------------------------------
# Page setup + styles
# ---------------------------------------------------------
st.set_page_config(
    page_title="TechTitans Renewable Valuation",
    page_icon="⚡",
    layout="wide",
)

st.markdown(
    """
    <style>
      .stApp { background-color: #f9fafc; }
      .main-title {
        font-size: 40px; font-weight: 800; color: #003366; margin-bottom: 0px;
        text-align: center;
      }
      .subtitle {
        font-size: 18px; color: #5a6b7a; margin-top: 6px; text-align: center;
      }
      .section-title { font-weight:700; color:#003366; }
      .caption-small { color:#6b7f90; font-size:13px; }
      .kpi { background:#fff; padding:18px; border-radius:14px; border:1px solid #edf0f3; }
    </style>
    """,
    unsafe_allow_html=True,
)

st.markdown("<div class='main-title'>🌍 TechTitans Renewable Valuation Dashboard</div>", unsafe_allow_html=True)
st.markdown(
    "<div class='subtitle'>Dynamic market risk & pricing tool • ERCOT • MISO • CAISO</div>",
    unsafe_allow_html=True,
)
st.markdown("")

with st.expander("🧭 Quick guide — how to use this dashboard", expanded=False):
    st.markdown("""
    **1) Choose Market & Project** in the left sidebar (ERCOT, MISO, CAISO).  
    **2) Set Risk Tolerance (P-level)** — higher P = more conservative (per prompt, **P75 ⇒ use the bottom 25%** of merchant outcomes).  
    **3) Choose Negative-Price Handling** — include negative prices or treat them as 0 revenue.  
    **4) Read the KPIs, Decision Brief, and the component table** (Hub Forward, Basis, Risk Adj).  
    **5) Explore visuals** in *Risk Comparison* and *Charts & Export*.  
    **6) Export** CSV/Excel/PPT for your submission.
    """)

# ---------------------------------------------------------
# Data access
# ---------------------------------------------------------
DATA_PATH = Path(__file__).resolve().parents[1] / "data" / "HackathonDataset.xlsx"
AVAILABLE_SHEETS = ["ERCOT", "MISO", "CAISO"]

@st.cache_data(show_spinner=False)
def load_sheet_raw(sheet: str) -> pd.DataFrame:
    """Load raw sheet with unknown header positions and normalize columns."""
    if not DATA_PATH.exists():
        raise FileNotFoundError(f"Missing data file: {DATA_PATH}")
    probe = pd.read_excel(DATA_PATH, sheet_name=sheet, header=None, nrows=30)
    header_row = None
    for i in range(len(probe)):
        row_vals = probe.iloc[i].astype(str).str.strip().str.lower().tolist()
        if "date" in row_vals and "he" in row_vals:
            header_row = i
            break
    if header_row is None:
        for i in (8, 9, 10, 7):  # common header rows
            if i < len(probe):
                row_vals = probe.iloc[i].astype(str).str.strip().str.lower().tolist()
                if "date" in row_vals and "he" in row_vals:
                    header_row = i
                    break
    if header_row is None:
        header_row = 0
    df = pd.read_excel(DATA_PATH, sheet_name=sheet, header=header_row).copy()
    df.columns = [str(c).strip() for c in df.columns]
    return df

def _to_num(series: pd.Series) -> pd.Series:
    """Convert '$' strings and '(x)' negatives to float, keep NaN for '-'."""
    if series is None:
        return pd.Series(dtype=float)
    ser = series.astype(str).str.replace(r"[\$,]", "", regex=True)
    ser = ser.str.replace(r"^\((.*)\)$", r"-\1", regex=True)  # (3.20) -> -3.20
    ser = ser.replace({"-": np.nan, "None": np.nan, "nan": np.nan})
    return pd.to_numeric(ser, errors="coerce")

@st.cache_data(show_spinner=False)
def load_standardized(sheet: str) -> pd.DataFrame:
    """
    Return standardized dataframe with:
    ['timestamp','gen','da_hub','da_node','rt_hub','rt_node']
    """
    df = load_sheet_raw(sheet)
    if "Date" not in df.columns or "HE" not in df.columns:
        raise KeyError(f"Expected 'Date' and 'HE' in {sheet}. Got: {list(df.columns)}")

    df["timestamp"] = pd.to_datetime(df["Date"], errors="coerce") + pd.to_timedelta(
        pd.to_numeric(df["HE"], errors="coerce") - 1, unit="h"
    )

    # Generation
    gen_col = "Gen" if "Gen" in df.columns else None
    df["gen"] = _to_num(df.get(gen_col, pd.Series(index=df.index, dtype="float")))

    # Price columns: ERCOT has RT Hub/Busbar; others use Hub/Busbar for RT
    if sheet == "ERCOT":
        df["da_hub"]  = _to_num(df.get("DA Hub"))
        df["da_node"] = _to_num(df.get("DA Busbar"))
        df["rt_hub"]  = _to_num(df.get("RT Hub"))
        df["rt_node"] = _to_num(df.get("RT Busbar"))
    else:
        df["da_hub"]  = _to_num(df.get("DA Hub"))
        df["da_node"] = _to_num(df.get("DA Busbar"))
        df["rt_hub"]  = _to_num(df.get("Hub"))
        df["rt_node"] = _to_num(df.get("Busbar"))

    std = df[["timestamp", "gen", "da_hub", "da_node", "rt_hub", "rt_node"]].copy()
    std = std.dropna(subset=["timestamp"]).sort_values("timestamp")
    return std

# ---------------------------------------------------------
# Core pricing logic
# ---------------------------------------------------------
def apply_negative_price_policy(df: pd.DataFrame, include_negatives: bool) -> pd.DataFrame:
    """If include_negatives is False => set price<0 to 0 (revenue=0) for all price columns."""
    out = df.copy()
    if not include_negatives:
        for c in ["da_hub", "da_node", "rt_hub", "rt_node"]:
            out[c] = np.where(out[c] < 0, 0.0, out[c])
    return out

def p_level(series: pd.Series, p: int) -> float:
    """
    Return P-level (quantile), consistent with hackathon clarification:
    P75 appetite ⇒ pick the bottom 25% (q=0.25) of merchant outcomes.
    """
    q = 1 - (p / 100.0)  # P75 -> 0.25
    s = series.dropna()
    return float(np.nanquantile(s, q=q)) if len(s) else np.nan

def compute_results(std: pd.DataFrame, p: int, include_negatives: bool) -> pd.DataFrame:
    """
    Compute P-level for DA/RT × Hub/Busbar; also decompose Hub Forward, Basis, and a protective Risk Adj.
    """
    df = apply_negative_price_policy(std, include_negatives)

    # Forward proxy & basis
    hub_fw_da = df["da_hub"].mean()
    basis_da = (df["da_node"] - df["da_hub"]).mean()
    hub_fw_rt = df["rt_hub"].mean()
    basis_rt = (df["rt_node"] - df["rt_hub"]).mean()

    # Protective risk adjustment: emphasize negative tail of basis (P10)
    tail_da = float(np.nanquantile((df["da_node"] - df["da_hub"]).dropna(), 0.10)) if df["da_node"].notna().any() else 0.0
    tail_rt = float(np.nanquantile((df["rt_node"] - df["rt_hub"]).dropna(), 0.10)) if df["rt_node"].notna().any() else 0.0
    risk_adj_da = min(0.0, tail_da) * 2.0
    risk_adj_rt = min(0.0, tail_rt) * 2.0

    rows = []
    variants = [
        ("DA Hub",    "da_hub",  hub_fw_da, basis_da, risk_adj_da),
        ("DA Busbar", "da_node", hub_fw_da, basis_da, risk_adj_da),
        ("RT Hub",    "rt_hub",  hub_fw_rt, basis_rt, risk_adj_rt),
        ("RT Busbar", "rt_node", hub_fw_rt, basis_rt, risk_adj_rt),
    ]
    for name, col, hub_forward, basis, risk_adj in variants:
        F = p_level(df[col], p)
        rows.append({
            "Price Variant": name,
            "P-Level ($/MWh)": round(F, 2) if pd.notna(F) else None,
            "Hub Forward": round(hub_forward, 2) if pd.notna(hub_forward) else None,
            "Basis": round(basis, 2) if pd.notna(basis) else None,
            "Risk Adj": round(risk_adj, 2) if pd.notna(risk_adj) else None
        })
    return pd.DataFrame(rows)

# ---------- EXPLANATION HELPERS ----------
def pct(x):
    try:
        return f"{100*float(x):.1f}%"
    except Exception:
        return "—"

def _safe_float(v, nd=2):
    try:
        return round(float(v), nd)
    except Exception:
        return None

def ai_insight_text(sheet: str, p: int, include_negatives: bool, results: pd.DataFrame) -> str:
    """Executive Decision Brief — concise recommendation + conservative target."""
    best = results.sort_values("P-Level ($/MWh)", ascending=False).iloc[0]
    policy = "including negative-price hours" if include_negatives else "excluding negative-price hours"
    target = None
    try:
        target = float(best["P-Level ($/MWh)"]) + float(best["Risk Adj"] or 0)
    except Exception:
        pass

    lines = []
    lines.append(f"**Market**: `{sheet}` • **Risk appetite**: **P{p}** • **Policy**: **{policy}**")
    lines.append(f"**Recommended product**: **{best['Price Variant']}**")
    lines.append(
        f"**Fixed-price anchor (P-level)**: **${best['P-Level ($/MWh)']}/MWh**  "
        f"**Hub fwd**: **${best['Hub Forward']}/MWh**, **Basis (avg)**: **${best['Basis']}/MWh**, "
        f"**Risk adj**: **${best['Risk Adj']}/MWh**"
    )
    if target is not None:
        lines.append(f"**Conservative target price**: **${round(target, 2)}/MWh** (P-level + risk adjustment)")
    lines.append(
        "_Why this matters_: P-level follows the prompt clarification — **P75 appetite ⇒ choose the bottom "
        "**25%** of the merchant distribution to protect downside. Basis & risk adjustment de-risk hub→node exposure."
    )
    return "\n\n".join(lines)

def market_commentary(sheet: str, p: int, include_negatives: bool, std_df: pd.DataFrame, results: pd.DataFrame) -> str:
    """Narrative under single-market tables/plots."""
    if std_df is None or std_df.empty or results is None or results.empty:
        return "_No data available for commentary._"

    neg_mask = (std_df[['da_hub', 'da_node', 'rt_hub', 'rt_node']] < 0)
    neg_share_any = neg_mask.any(axis=1).mean()
    neg_msg = "including" if include_negatives else "excluding (counted as $0 revenue)"

    if 'da_node' in std_df and 'da_hub' in std_df:
        basis_series = std_df['da_node'] - std_df['da_hub']
    else:
        basis_series = std_df['rt_node'] - std_df['rt_hub']
    b_avg = _safe_float(basis_series.mean())
    b_p10 = _safe_float(basis_series.quantile(0.10))
    b_p90 = _safe_float(basis_series.quantile(0.90))

    best = results.sort_values("P-Level ($/MWh)", ascending=False).iloc[0]
    best_name = str(best['Price Variant'])
    best_p = _safe_float(best["P-Level ($/MWh)"])
    best_basis = _safe_float(best["Basis"])
    best_rf = _safe_float(best["Risk Adj"])

    stance = {
        50: "balanced (median outcome)",
        60: "mildly conservative",
        70: "conservative",
        75: "downside-protective (hackathon default)",
        80: "very conservative",
        85: "highly conservative",
        90: "strong downside protection"
    }
    stance_text = stance.get(int(p), "custom")

    lines = []
    lines.append(f"**Risk view**: P**{p}** ({stance_text}) with negative-price policy **{neg_msg}**.")
    lines.append(
        f"**Basis behavior (node − hub)**: avg **${b_avg}/MWh**, "
        f"P10 **${b_p10}/MWh**, P90 **${b_p90}/MWh**."
    )
    lines.append(f"**Negative-price exposure**: {pct(neg_share_any)} of hours had at least one negative price.")
    lines.append(
        f"**Recommended variant**: **{best_name}** with P-level **${best_p}/MWh**; "
        f"components: basis ≈ **${best_basis}/MWh**, risk adj ≈ **${best_rf}/MWh**."
    )
    if p >= 85:
        lines.append("At this stance, we prioritize **downside tails**; expect lower targets and stronger protective discounts.")
    elif p >= 75:
        lines.append("This stance aims to be **robust to adverse months**; prices balance forwards with basis downside.")
    else:
        lines.append("Lower P-levels accept more **merchant variability**; targets sit closer to plain forwards.")

    if (b_p10 is not None) and b_p10 < 0:
        lines.append("**Takeaway**: Tail basis shows **node discount vs hub**; hub-only hedges can underperform without basis protection.")
    elif (b_p10 is not None) and b_p10 > 0:
        lines.append("**Takeaway**: Tail basis is **positive**; hub-settled hedges may leave node upside/variability unhedged.")
    else:
        lines.append("**Takeaway**: Basis tails are muted; hub vs node choice is more operational than risk-driven.")

    return "\n\n".join(lines)

def comparison_commentary(selected_markets: list, p: int, include_negatives: bool, compare_df: pd.DataFrame) -> str:
    """Narrative under the multi-market comparison chart."""
    if compare_df is None or compare_df.empty:
        return "_No comparison data to summarize._"
    row_max = compare_df.loc[compare_df['P-Level ($/MWh)'].idxmax()]
    row_min = compare_df.loc[compare_df['P-Level ($/MWh)'].idxmin()]
    max_msg = f"Highest P-level: **{row_max['Market']} – {row_max['Price Variant']}** at **${_safe_float(row_max['P-Level ($/MWh)'])}/MWh**"
    min_msg = f"Lowest P-level: **{row_min['Market']} – {row_min['Price Variant']}** at **${_safe_float(row_min['P-Level ($/MWh)'])}/MWh**"
    lines = []
    lines.append(f"**Risk appetite**: P**{p}**; negative-price policy is "
                 f"{'included' if include_negatives else 'excluded (0 revenue)'} for all calculations shown.")
    lines.append(max_msg + " • " + min_msg)
    try:
        basis_sorted = compare_df[['Market','Price Variant','Basis']].sort_values('Basis')
        head = basis_sorted.head(1).iloc[0]
        tail = basis_sorted.tail(1).iloc[0]
        lines.append(f"**Most negative basis**: **{head['Market']} – {head['Price Variant']}** at **${_safe_float(head['Basis'])}/MWh**.")
        lines.append(f"**Most positive basis**: **{tail['Market']} – {tail['Price Variant']}** at **${_safe_float(tail['Basis'])}/MWh**.")
    except Exception:
        pass
    lines.append("**Interpretation**: Negative basis favors **busbar/node-settled** strategies; positive basis can make **hub-settled** look richer but leaves node risk.")
    return "\n\n".join(lines)

# ---------------------------------------------------------
# Sidebar: navigation and inputs
# ---------------------------------------------------------
st.sidebar.markdown("### Navigation")
menu = st.sidebar.radio(
    "Go to",
    ["🏠 Overview", "📊 Market Analysis", "📉 Risk Comparison", "📈 Charts & Export"],
)

st.sidebar.markdown("---")
st.sidebar.markdown("### Inputs")
market_label = st.sidebar.selectbox(
    "Select Market & Project",
    ["ERCOT – Valentino (Wind)", "MISO – Mantero (Wind)", "CAISO – Howling Gale (Solar)"],
)
sheet_key = "ERCOT" if market_label.startswith("ERCOT") else ("MISO" if market_label.startswith("MISO") else "CAISO")

p_level_input = st.sidebar.slider("Choose Risk Tolerance (P-level)", min_value=50, max_value=95, value=75, step=5)
neg_choice = st.sidebar.radio(
    "Negative Price Handling",
    ["Include negative-price hours", "Exclude (revenue=0 when price < 0)"],
    index=0,
)
include_negative = (neg_choice == "Include negative-price hours")

with st.sidebar.expander("ℹ️ Tips"):
    st.write(
        "• P75 is the hackathon default (uses the **bottom 25%** of merchant outcomes).\n"
        "• Excluding negatives simulates PPAs that do not pay when price < 0."
    )

# ---------------------------------------------------------
# Load data + compute
# ---------------------------------------------------------
try:
    std = load_standardized(sheet_key)
except Exception as e:
    st.error(f"Failed to load {sheet_key}: {e}")
    st.stop()

results = compute_results(std, p_level_input, include_negative)

# KPI helpers
kpi_map = {r["Price Variant"]: r for _, r in results.iterrows()}
da_hub_kpi = kpi_map.get("DA Hub", {}).get("P-Level ($/MWh)", None)
rt_hub_kpi = kpi_map.get("RT Hub", {}).get("P-Level ($/MWh)", None)
avg_basis_kpi = round(results["Basis"].mean(), 2)

# ---------------------------------------------------------
# Overview
# ---------------------------------------------------------
if menu == "🏠 Overview":
    st.markdown("### Executive Summary")

    c1, c2, c3 = st.columns([1, 1, 1])
    with c1:
        st.markdown("<div class='kpi'>", unsafe_allow_html=True)
        st.metric("DA Hub P-Level", f"${da_hub_kpi}")
        st.markdown("</div>", unsafe_allow_html=True)
    with c2:
        st.markdown("<div class='kpi'>", unsafe_allow_html=True)
        st.metric("RT Hub P-Level", f"${rt_hub_kpi}")
        st.markdown("</div>", unsafe_allow_html=True)
    with c3:
        st.markdown("<div class='kpi'>", unsafe_allow_html=True)
        st.metric("Average Basis", f"${avg_basis_kpi}")
        st.markdown("</div>", unsafe_allow_html=True)

    st.markdown("---")
    st.markdown("#### Decision Brief")
    st.success(ai_insight_text(sheet_key, p_level_input, include_negative, results))

    st.markdown("---")
    with st.expander("Preview a sample of input data (first 20 rows)"):
        st.dataframe(std.head(20), use_container_width=True)
        st.caption("Data source: `data/HackathonDataset.xlsx`")

    with st.expander("📚 Glossary"):
        st.markdown("""
        - **P-level**: Downside-focused anchor. *P75 appetite* ⇒ take **P25** of merchant distribution.  
        - **Hub Forward**: Average DA hub price as a simple forward proxy.  
        - **Basis**: Node − Hub (avg). Positive = node richer than hub; negative = discount.  
        - **Risk Adjustment**: Protective deduction from basis tail risk (10th percentile).  
        - **Negative-price policy**: If OFF, hours with price < 0 are treated as $0 revenue.  
        """)

# ---------------------------------------------------------
# Market Analysis (table + commentary + downloads)
# ---------------------------------------------------------
elif menu == "📊 Market Analysis":
    st.markdown("### Summary of Computed P-Level & Risk Metrics")
    st.dataframe(results, use_container_width=True)
    st.caption("P-level uses downside quantiles of the merchant outcome distribution. Per prompt: **P75 ⇒ P25**.")

    # Commentary under table
    st.markdown("##### What this means")
    st.markdown(
        market_commentary(
            sheet=sheet_key,
            p=int(p_level_input),
            include_negatives=include_negative,
            std_df=std,
            results=results
        )
    )

    # Downloads
    st.markdown("---")
    csv_bytes = results.to_csv(index=False).encode("utf-8")
    st.download_button("⬇️ Download results (CSV)", csv_bytes, file_name=f"{sheet_key}_results.csv", mime="text/csv")

    excel_buf = io.BytesIO()
    with pd.ExcelWriter(excel_buf, engine="xlsxwriter") as writer:
        results.to_excel(writer, index=False, sheet_name="Results")
        std.head(2000).to_excel(writer, index=False, sheet_name="DataPreview")
    st.download_button(
        "⬇️ Download results (Excel)",
        excel_buf.getvalue(),
        file_name=f"{sheet_key}_results.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )

# ---------------------------------------------------------
# Risk Comparison (multi-market)
# ---------------------------------------------------------
elif menu == "📉 Risk Comparison":
    st.markdown("### Compare Markets")

    selected_markets = st.multiselect(
        "Pick 1–3 markets to compare",
        ["ERCOT – Valentino (Wind)", "MISO – Mantero (Wind)", "CAISO – Howling Gale (Solar)"],
        default=[market_label],
        max_selections=3,
    )

    comp_rows = []
    for label in selected_markets:
        key = "ERCOT" if label.startswith("ERCOT") else ("MISO" if label.startswith("MISO") else "CAISO")
        df_std = load_standardized(key)
        df_res = compute_results(df_std, p_level_input, include_negative)
        df_res["Market"] = label
        comp_rows.append(df_res)

    if comp_rows:
        compare_table = pd.concat(comp_rows, ignore_index=True)
        st.dataframe(compare_table, use_container_width=True)

        st.markdown("#### P-Level by Market and Variant")
        fig_bar = px.bar(
            compare_table,
            x="Market",
            y="P-Level ($/MWh)",
            color="Price Variant",
            barmode="group",
            text="P-Level ($/MWh)",
            color_discrete_sequence=px.colors.qualitative.Set2,
            height=420,
        )
        fig_bar.update_traces(texttemplate="%{text}", textposition="outside")
        fig_bar.update_layout(margin=dict(l=10, r=10, t=30, b=10))
        st.plotly_chart(fig_bar, use_container_width=True)
        st.caption("Note: Bars reflect current P-level and negative-price policy choices.")

        # Cross-market interpretation
        st.markdown("##### Cross-market interpretation")
        st.markdown(
            comparison_commentary(
                selected_markets=selected_markets,
                p=int(p_level_input),
                include_negatives=include_negative,
                compare_df=compare_table
            )
        )

# ---------------------------------------------------------
# Charts & Export (scatter + PPT)
# ---------------------------------------------------------
elif menu == "📈 Charts & Export":
    st.markdown("### Risk-Adjusted Pricing vs Market Forward (Scatter)")

    scatter_df = results.copy()
    fig = px.scatter(
        scatter_df,
        x="Hub Forward",
        y="P-Level ($/MWh)",
        color="Price Variant",
        size=np.clip(scatter_df["P-Level ($/MWh)"], 12, 28),
        text="Price Variant",
        color_discrete_sequence=px.colors.qualitative.Set1,
        height=420,
    )
    fig.update_traces(textposition="top center")
    fig.update_layout(xaxis_title="Hub Forward ($/MWh)", yaxis_title="P-Level ($/MWh)")
    st.plotly_chart(fig, use_container_width=True)
    st.caption("Note: P-level anchors the downside percentile of merchant outcomes; points reflect your current policy & basis adjustments.")

    # Single-market narrative under scatter
    st.markdown("##### What this scatter suggests")
    st.markdown(
        market_commentary(
            sheet=sheet_key,
            p=int(p_level_input),
            include_negatives=include_negative,
            std_df=std,
            results=results
        )
    )

    st.markdown("---")
    st.markdown("### Export Report")

    def make_ppt(rst: pd.DataFrame, market: str, p: int, include_neg: bool) -> bytes:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title-only layout

        # Title
        title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(11), Inches(1.0))
        tf = title.text_frame
        p1 = tf.paragraphs[0]
        p1.text = "TechTitans – Renewable Valuation Summary"
        p1.font.size = Pt(28)
        p1.font.bold = True

        # Subtitle
        sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(11), Inches(0.8))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        policy = "Include negatives" if include_neg else "Exclude negatives"
        p2.text = f"Market: {market} | P-level: P{p} | Policy: {policy}"
        p2.font.size = Pt(16)

        # Table
        rows, cols = rst.shape[0] + 1, rst.shape[1]
        table_shape = slide.shapes.add_table(rows, cols, Inches(0.7), Inches(2.0), Inches(11), Inches(1.0 + 0.3*rows))
        table = table_shape.table
        for j, col in enumerate(rst.columns):
            table.cell(0, j).text = str(col)
        for i in range(rst.shape[0]):
            for j in range(cols):
                table.cell(i + 1, j).text = str(rst.iloc[i, j])

        bio = io.BytesIO()
        prs.save(bio)
        return bio.getvalue()

    ppt_bytes = make_ppt(results, sheet_key, p_level_input, include_negative)
    st.download_button(
        "📊 Download PowerPoint summary",
        ppt_bytes,
        file_name=f"{sheet_key}_pricing_summary.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

# ---------------------------------------------------------
# Footer
# ---------------------------------------------------------
st.markdown("<hr>", unsafe_allow_html=True)
st.markdown(
    "<p style='text-align:center; color:#6b7f90'>Built with ❤️ by <b>TechTitans</b> • Avangrid Hackathon 2025 • "
    "Powered by Streamlit, Pandas & Plotly</p>",
    unsafe_allow_html=True,
)
